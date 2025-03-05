from pymongo import MongoClient
from pymongo.collection import Collection
from pymongo.server_api import ServerApi
import PyPDF2
from langchain_community.document_loaders import PyPDFLoader

from pptx import Presentation
import re

def extract_questions_and_answers(pdf_path):
    questions_answers = []
    
    with open(f"./{pdf_path}", "rb") as file:
        reader = PyPDF2.PdfReader(file)
        text = "".join([page.extract_text() for page in reader.pages if page.extract_text()])
    
    # Expresión regular para capturar preguntas numeradas y sus respuestas
    pattern = re.compile(r"(\d+)\.\s?(.*?)\?\s*(.*?)(?=\n\d+\.| \Z)", re.S)
    
    matches = pattern.findall(text)
    
    for match in matches:
        question_number, question_text, answer_text = match
        question_text = question_text.strip()
        question_text = question_text.replace("\n", "")

        answer_text = answer_text.strip()
        answer_text = answer_text.replace("\n", "")
        
        questions_answers.append({"question": question_text.strip(), "answer": answer_text})
    
    return questions_answers


def add_document_to_mongo(database, collection_name,pdf_path, embeddings_model):
    """
    Añade un documento a MongoDB con su respectivo embedding semántico.
    El documento debe de tener una estructura similar a la siguiente:
    
    1. Pregunta 1?
    Respuesta 1
    
    2. Pregunta 2?
    Respuesta 2
    
    ..."""
    collection = database[collection_name]

    datos = extract_questions_and_answers(pdf_path)
    print(datos)

    for data in datos:
        question = data["question"]
        answer = data["answer"]

        embedding = embeddings_model.encode(question)
        document = {"question": question, "answer":answer, "semantic_embedding": embedding.tolist()}
        print(f"Insertando documento: {document}")
        collection.insert_one(document)

    

def update_mongo_with_embeddings(database, collection_name, embeddings_model):
    """
    Actualiza los documentos en MongoDB añadiendo la propiedad 'semantic_embedding'.
    
    :param database: Objeto de la base de datos MongoDB.
    :param collection_name: Nombre de la colección en MongoDB.
    :param embeddings_model: Modelo que genera los embeddings.
    """
    collection = database[collection_name]
    documents = list(collection.find({}))
    
    for doc in documents:
        text = f"{doc.get('question', '')}"
        embedding = embeddings_model.encode(text)

        collection.update_one({"_id": doc["_id"]}, {"$set": {"semantic_embedding": embedding.tolist()}})

def search_mongodb(query:str, collection, semantic_search):
    """
    Busca documentos en MongoDB utilizando embeddings semánticos.
    
    :param query: Texto de la consulta.
    :return: Lista de documentos que coinciden con la consulta.
    """
    # collection = db["answer"]
    query_embedding = semantic_search.model.encode(query)
    #la siguiente consulta busca los 3 documentos más similares a la consulta.La 2 consulta es para recuperar los campos que se desean mostrar. En este caso, se muestra el campo de "answer"
    results = collection.aggregate([{
    "$vectorSearch": {
      "index": "vector_index",
      "path": "semantic_embedding",
      "queryVector": query_embedding.tolist(),
      "numCandidates": 3,
      "limit": 3
        }
    }, {
          "$project": {
             "_id": 0,
             "answer": 1,
             "score": {
                "$meta": "vectorSearchScore"
             }
          }
       }
  
    ])


    return list(results)


def read_pdf(file_path: str) -> str:
    """Lee un archivo PDF y devuelve su contenido formateado."""
    text = []
    
    try:
        with open(file_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text_page = page.extract_text()
                text_page = text_page.replace("\n", "")
                text.append(text_page)
                # text.append(page.extract_text() or "")  # Extraer texto de cada página
    except Exception as e:
        return f"Error al leer el PDF: {str(e)}"
    
    return "\n".join(text).strip()
    # return text

# def extract_text_from_powerpoint(file_path):
#     prs = Presentation(file_path)
#     text = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 text += shape.text_frame.text
#     return text

# def extract_text_from_other_file_types(file_path):
#     with open(file_path, 'r', encoding='utf-8') as file:
#         text = file.read()
#     return text
    
def conexionMongoDB():
    try:
        uri = "mongodb+srv://chatbot:Rufo2009@cluster0.yju2g.mongodb.net/?retryWrites=true&w=majority&appName=Cluster0"
        # Create a new client and connect to the server
        client = MongoClient(uri, server_api=ServerApi('1'), port=27017)
        database = client['chatbot-embbeding']
    except Exception as ex:
        print("Error durante la conexión: {}".format(ex))
    return database

# database = conexionMongoDB()

# semantic_search = SemanticSearchEnhancer()  

# add_document_to_mongo(database,collection_name= "answer", pdf_path="Preguntas_cuestionario_resumidas.pdf", embeddings_model=semantic_search.model)

# busqueda = search_mongodb("¿Que soluciones ofrece Mobbeel?", database, semantic_search)

# print(list(busqueda))
texto = read_pdf("Preguntas_cuestionario_resumidas.pdf")
print(texto)